"""Build Whitespace hackathon deck — RedClaw brand colours."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── RedClaw palette ──────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x02, 0x23, 0x66)
REDCLAW_BLUE  = RGBColor(0x0B, 0x41, 0xCD)
LIGHT_BLUE  = RGBColor(0x14, 0x82, 0xFA)
XL_BLUE     = RGBColor(0xBD, 0xE3, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF5, 0xF5, 0xF2)
GREY3       = RGBColor(0xC2, 0xBA, 0xB5)
GREY1       = RGBColor(0x54, 0x4F, 0x4F)
GREEN       = RGBColor(0x00, 0xB4, 0x58)
ORANGE      = RGBColor(0xFF, 0x7D, 0x29)
RED         = RGBColor(0xFF, 0x1F, 0x26)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # truly blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=14, bold=False, color=GREY1,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = _Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def bullet(tf, text, size=13, color=GREY1, indent=0, bold=False):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = indent
    run = p.add_run()
    run.text = "• " + text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    return p


def header_bar(slide, title, subtitle=None):
    """Dark-blue top bar with white title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill=DARK_BLUE)
    add_text(slide, title,
             Inches(0.45), Inches(0.12), Inches(10), Inches(0.65),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.45), Inches(0.72), Inches(10), Inches(0.38),
                 size=14, color=XL_BLUE)
    # RedClaw blue accent line
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Pt(3), fill=REDCLAW_BLUE)


def footer(slide, label="Whitespace · RedClaw CSI Hackathon 2026 · Confidential"):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), fill=DARK_BLUE)
    add_text(slide, label,
             Inches(0.3), Inches(7.17), Inches(10), Inches(0.3),
             size=9, color=GREY3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Full background
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BLUE)
# RedClaw-blue diagonal accent band
add_rect(sl, Inches(8.5), 0, Inches(4.83), SLIDE_H, fill=REDCLAW_BLUE)
# Light overlay strip
add_rect(sl, Inches(8.3), 0, Inches(0.25), SLIDE_H, fill=LIGHT_BLUE)

add_text(sl, "Whitespace",
         Inches(0.6), Inches(1.8), Inches(7.8), Inches(1.4),
         size=54, bold=True, color=WHITE)

add_text(sl, "Autonomous Drug Discovery Intelligence",
         Inches(0.6), Inches(3.1), Inches(7.8), Inches(0.7),
         size=22, color=XL_BLUE)

add_text(sl, "From fragmented databases to CEO-ready strategy — in one agent run.",
         Inches(0.6), Inches(3.75), Inches(7.6), Inches(0.65),
         size=15, color=GREY3, italic=True)

add_text(sl, "RedClaw CSI Hackathon 2026",
         Inches(0.6), Inches(5.8), Inches(5), Inches(0.4),
         size=12, color=GREY3)
add_text(sl, "April 2026  |  For Internal Use Only",
         Inches(0.6), Inches(6.15), Inches(5), Inches(0.35),
         size=11, color=GREY3)

# Right-panel label
add_text(sl, "Strategic\nIntelligence\nAgent",
         Inches(8.85), Inches(2.5), Inches(4), Inches(2.5),
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "34 tools  ·  GPU-accelerated\nReAct loop  ·  Audit trail",
         Inches(8.85), Inches(5.0), Inches(4), Inches(1.0),
         size=13, color=XL_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "The Problem", "Strategic intelligence is siloed, slow, and manual")
footer(sl)

# Quote callout
add_rect(sl, Inches(0.4), Inches(1.4), Inches(12.5), Inches(1.2), fill=XL_BLUE)
add_text(sl, '"Where should RedClaw move in the next 90 days?"',
         Inches(0.65), Inches(1.5), Inches(12), Inches(0.9),
         size=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, italic=True)

# Pain points — 3 columns
cols = [
    ("Days of analyst work", "Cross-referencing ClinicalTrials.gov, Open Targets, ChEMBL, Europe PMC, ArXiv, USPTO, FDA FAERS, UniProt, KEGG, Orphanet manually."),
    ("Insights arrive stale", "By the time a briefing lands, competitors have moved. A three-week-old sweep misses fast-moving clinical updates."),
    ("No audit trail", "Manual analysis cannot be reproduced, validated, or scaled. One analyst, one perspective, no consistency."),
]
x = Inches(0.4)
for title, body in cols:
    add_rect(sl, x, Inches(2.85), Inches(3.95), Inches(3.85), fill=WHITE, line=REDCLAW_BLUE)
    add_text(sl, title, x + Inches(0.15), Inches(3.0), Inches(3.65), Inches(0.55),
             size=15, bold=True, color=DARK_BLUE)
    txb = sl.shapes.add_textbox(x + Inches(0.15), Inches(3.55), Inches(3.65), Inches(2.8))
    tf = txb.text_frame
    tf.word_wrap = True
    add_para(tf, body, size=13, color=GREY1)
    x += Inches(4.3)

# Bottom tag
add_text(sl, "This is an agent-shaped problem: multi-step · tool-heavy · reasoning-intensive · time-sensitive",
         Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.35),
         size=12, bold=True, color=REDCLAW_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "Solution: Whitespace", "A 34-tool ReAct agent for pharmaceutical strategic intelligence")
footer(sl)

# Left: what it does
add_rect(sl, Inches(0.4), Inches(1.35), Inches(5.9), Inches(5.45), fill=WHITE, line=REDCLAW_BLUE)
add_text(sl, "What it does in one job", Inches(0.6), Inches(1.45), Inches(5.5), Inches(0.45),
         size=15, bold=True, color=DARK_BLUE)
txb = sl.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(5.5), Inches(4.6))
tf = txb.text_frame; tf.word_wrap = True
for b in [
    "Ranks RedClaw's 59 pipeline assets by composite opportunity score",
    "Identifies competitive white spaces vs 8 major pharma companies",
    "Scores pipeline assets for trial success probability",
    "Flags competitive threats & patent exposures",
    "Generates branded PDF report — unattended",
    "Logs every tool call to a JSON Lines audit trail",
]:
    bullet(tf, b, size=13, color=GREY1)

# Right: flow diagram (text-based)
add_rect(sl, Inches(6.6), Inches(1.35), Inches(6.3), Inches(5.45), fill=DARK_BLUE)
add_text(sl, "Agent Loop", Inches(6.8), Inches(1.45), Inches(5.9), Inches(0.45),
         size=15, bold=True, color=WHITE)

steps = [
    (REDCLAW_BLUE, "Claude claude-opus-4-6  (ReAct reasoning)"),
    (LIGHT_BLUE, "Tool selection & sequencing  (34 tools)"),
    (LIGHT_BLUE, "Real API calls  (trials, literature, patents, genes)"),
    (LIGHT_BLUE, "GPU inference  (fold · variant · ADMET · dock)"),
    (LIGHT_BLUE, "Stall detection & recovery"),
    (GREEN,      "PDF report  +  audit trail"),
]
y = Inches(2.0)
for col, label in steps:
    add_rect(sl, Inches(6.8), y, Inches(5.9), Inches(0.62), fill=col)
    add_text(sl, label, Inches(6.95), y + Inches(0.08), Inches(5.6), Inches(0.5),
             size=12, bold=(col == GREEN), color=WHITE if col != GREEN else DARK_BLUE)
    y += Inches(0.72)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TOOL LANDSCAPE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "Tool Landscape", "34 tools across 7 domains — all callable by the agent")
footer(sl)

categories = [
    ("Discovery",        DARK_BLUE,  ["find_gaps", "get_biology", "search_redclaw_trials", "find_combinations", "get_pathway_context (KEGG)"]),
    ("Competitive",      REDCLAW_BLUE, ["rank_portfolio", "list_pipeline_assets", "map_regulatory_path", "score_trial_outcome", "check_orphan_eligibility", "get_disease_prevalence"]),
    ("Literature",       LIGHT_BLUE, ["scan_literature", "scan_arxiv", "bulk_scan_literature"]),
    ("Chemistry",        RGBColor(0xFF,0x7D,0x29), ["find_hits + SMILES descriptors", "find_repurposing_candidates", "query_adverse_events (FDA FAERS)"]),
    ("GPU / Structure",  GREEN,      ["fold_target (Boltz-1)", "score_variant_effect (ESM-2)", "predict_admet", "cluster_scaffolds", "dock_compound"]),
    ("Patents",          GREY1,      ["search_patents (USPTO)", "get_patent_landscape (Lens.org)"]),
    ("Memory",           RGBColor(0xBC,0x36,0xF0), ["recall_longterm_memory", "save_to_cache", "generate_pdf_report"]),
]

x, y = Inches(0.3), Inches(1.35)
col_w, row_h = Inches(1.82), Inches(0.0)

positions = [
    (Inches(0.3),  Inches(1.35), Inches(3.8),  Inches(2.05)),
    (Inches(4.3),  Inches(1.35), Inches(3.8),  Inches(2.55)),
    (Inches(8.3),  Inches(1.35), Inches(4.6),  Inches(1.45)),
    (Inches(0.3),  Inches(3.6),  Inches(3.8),  Inches(1.95)),
    (Inches(4.3),  Inches(4.1),  Inches(3.8),  Inches(2.45)),
    (Inches(8.3),  Inches(3.0),  Inches(4.6),  Inches(1.25)),
    (Inches(8.3),  Inches(4.45), Inches(4.6),  Inches(1.5)),
]

for i, (cat, col, tools) in enumerate(categories):
    lx, ly, lw, lh = positions[i]
    add_rect(sl, lx, ly, lw, lh, fill=WHITE, line=col)
    # header strip
    add_rect(sl, lx, ly, lw, Inches(0.38), fill=col)
    add_text(sl, cat, lx + Inches(0.1), ly + Inches(0.03), lw - Inches(0.2), Inches(0.34),
             size=12, bold=True, color=WHITE)
    txb = sl.shapes.add_textbox(lx + Inches(0.12), ly + Inches(0.42), lw - Inches(0.2), lh - Inches(0.5))
    tf = txb.text_frame; tf.word_wrap = True
    for t in tools:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "· " + t
        run.font.size = Pt(11)
        run.font.color.rgb = GREY1


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — GPU PIPELINE (GenomeClaw)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "GPU Pipeline: GenomeClaw", "Boltz-1 folding · ESM-2 variant effects · ADMET · Docking — all on A100")
footer(sl)

# Pipeline flow
steps_gpu = [
    ("find_hits", "ChEMBL\nhit list\n+ SMILES", REDCLAW_BLUE),
    ("cluster_scaffolds", "Tanimoto\nscaffold\nclusters", REDCLAW_BLUE),
    ("fold_target", "Boltz-1\nprotein\nstructure", DARK_BLUE),
    ("dock_compound", "Geometry\ndocking\nscores", DARK_BLUE),
    ("predict_admet", "ADMET\nTIER-1\ngate ✓", GREEN),
    ("score_variant_effect", "ESM-2\nvariant\nΔΔG", DARK_BLUE),
]

x = Inches(0.35)
for name, label, col in steps_gpu:
    add_rect(sl, x, Inches(2.0), Inches(1.95), Inches(2.2), fill=col)
    add_text(sl, name, x + Inches(0.08), Inches(2.08), Inches(1.8), Inches(0.5),
             size=11, bold=True, color=WHITE)
    add_text(sl, label, x + Inches(0.08), Inches(2.55), Inches(1.8), Inches(1.4),
             size=12, color=WHITE, align=PP_ALIGN.CENTER)
    if x < Inches(10.5):
        add_text(sl, "→", x + Inches(1.95), Inches(2.85), Inches(0.4), Inches(0.6),
                 size=22, bold=True, color=REDCLAW_BLUE)
    x += Inches(2.28)

# Rust API box
add_rect(sl, Inches(0.35), Inches(4.5), Inches(12.65), Inches(0.65), fill=DARK_BLUE)
add_text(sl, "GenomeClaw Rust API  (clawapi)  —  served locally at 127.0.0.1:CLAWAPI_PORT  —  OpenSSL 3 via conda, WGPU Vulkan backend",
         Inches(0.5), Inches(4.57), Inches(12.3), Inches(0.5),
         size=11, color=XL_BLUE, align=PP_ALIGN.CENTER)

# Weights
add_rect(sl, Inches(0.35), Inches(5.35), Inches(6.1), Inches(1.3), fill=WHITE, line=REDCLAW_BLUE)
add_text(sl, "Weights", Inches(0.55), Inches(5.42), Inches(5.7), Inches(0.35),
         size=13, bold=True, color=DARK_BLUE)
txb = sl.shapes.add_textbox(Inches(0.55), Inches(5.8), Inches(5.7), Inches(0.75))
tf = txb.text_frame; tf.word_wrap = True
for w in ["boltz-1  (boltz1.ckpt — ~1.7GB, SLURM env)", "esm2_t33_650M_UR50D  (model.safetensors — 1.2GB)"]:
    bullet(tf, w, size=12, color=GREY1)

# Compute
add_rect(sl, Inches(6.65), Inches(5.35), Inches(6.3), Inches(1.3), fill=WHITE, line=DARK_BLUE)
add_text(sl, "Compute", Inches(6.85), Inches(5.42), Inches(5.9), Inches(0.35),
         size=13, bold=True, color=DARK_BLUE)
txb = sl.shapes.add_textbox(Inches(6.85), Inches(5.8), Inches(5.9), Inches(0.75))
tf = txb.text_frame; tf.word_wrap = True
for w in ["NVIDIA A100-SXM4-80GB on RedClaw sHPC", "SLURM batch_gpu / dia_gpu partition  ·  Singularity container"]:
    bullet(tf, w, size=12, color=GREY1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE & SECURITY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "Architecture & Security", "Designed for RedClaw sHPC — isolated, auditable, prompt-injection resistant")
footer(sl)

# Left column — arch
add_rect(sl, Inches(0.35), Inches(1.35), Inches(6.2), Inches(5.45), fill=WHITE, line=REDCLAW_BLUE)
add_text(sl, "Stack", Inches(0.55), Inches(1.45), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=DARK_BLUE)

layers = [
    (DARK_BLUE,  "Claude claude-opus-4-6  via ona-claude OAuth proxy"),
    (REDCLAW_BLUE, "proxy_server.py  —  retry loop (2× JSON-forcing preamble)"),
    (REDCLAW_BLUE, "run_agent.py  —  ReAct loop, stall counter, guard logic"),
    (LIGHT_BLUE, "34 Python tool functions  (requests → approved APIs)"),
    (LIGHT_BLUE, "GenomeClaw Rust API  (Boltz-1 · ESM-2 · ADMET · Docking)"),
    (GREEN,      "SLURM + Singularity  —  per-job isolated ports"),
]
y = Inches(1.95)
for col, label in layers:
    add_rect(sl, Inches(0.55), y, Inches(5.8), Inches(0.55), fill=col)
    add_text(sl, label, Inches(0.7), y + Inches(0.08), Inches(5.5), Inches(0.42),
             size=11, color=WHITE if col != GREEN else DARK_BLUE, bold=(col == GREEN))
    y += Inches(0.65)

# Right column — security
add_rect(sl, Inches(6.8), Inches(1.35), Inches(6.1), Inches(5.45), fill=WHITE, line=DARK_BLUE)
add_text(sl, "Security & Governance", Inches(7.0), Inches(1.45), Inches(5.7), Inches(0.4),
         size=15, bold=True, color=DARK_BLUE)

sec_items = [
    ("Egress Allowlist", "Monkey-patches requests.get/post at import time. Any outbound call to a non-approved domain raises PermissionError before TCP connects — defence against prompt-injection exfiltration."),
    ("Per-Job Port Isolation", "ONA_PORT = 8095 + (JOB_ID % 900). No two SLURM jobs on the same node can collide."),
    ("JSON Lines Audit Trail", "Every tool call logged: name, inputs digest, result digest, elapsed, SLURM job ID, container flag. Full reproducibility."),
    ("ADMET Hard Gate", "predict_admet is mandatory before any compound advances. TIER-2/3 compounds are dropped by the agent — no bypass path."),
    ("No PHI", "Only anonymised compound IDs, gene names, indication names. No patient-level data."),
]
y = Inches(2.0)
for title, body in sec_items:
    add_text(sl, title, Inches(7.0), y, Inches(5.7), Inches(0.35),
             size=12, bold=True, color=DARK_BLUE)
    add_text(sl, body, Inches(7.0), y + Inches(0.33), Inches(5.7), Inches(0.62),
             size=11, color=GREY1, wrap=True)
    y += Inches(1.02)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RESULTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "Results", "What Whitespace produced during hackathon week")
footer(sl)

# 4 big stat boxes
stats = [
    ("1",      "CEO Strategic\nBriefing",      "Analyst days → 45 min,\nunattended",        DARK_BLUE),
    ("8 × 2",  "Competitive +\nRepurposing",   "8-company sweeps\nin parallel overnight",   REDCLAW_BLUE),
    ("34",     "Tools called\nautonomously",   "Agent decides order\n& depth each run",      LIGHT_BLUE),
    ("85 %",   "Test coverage",                "314 tests · full audit\nmodule at 100%",      GREEN),
]
x = Inches(0.35)
for num, label, sub, col in stats:
    add_rect(sl, x, Inches(1.4), Inches(3.0), Inches(2.5), fill=col)
    add_text(sl, num, x + Inches(0.1), Inches(1.55), Inches(2.8), Inches(1.1),
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, label, x + Inches(0.1), Inches(2.6), Inches(2.8), Inches(0.65),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x, Inches(3.9), Inches(3.0), Inches(0.9), fill=WHITE, line=col)
    add_text(sl, sub, x + Inches(0.1), Inches(3.95), Inches(2.8), Inches(0.8),
             size=11, color=GREY1, align=PP_ALIGN.CENTER)
    x += Inches(3.27)

# GPU validation
add_rect(sl, Inches(0.35), Inches(5.05), Inches(12.65), Inches(1.6), fill=WHITE, line=REDCLAW_BLUE)
add_text(sl, "KRAS G12C Full Pipeline Validation",
         Inches(0.55), Inches(5.12), Inches(12.2), Inches(0.38),
         size=14, bold=True, color=DARK_BLUE)
txb = sl.shapes.add_textbox(Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.0))
tf = txb.text_frame; tf.word_wrap = True
flow = tf.add_paragraph()
flow.alignment = PP_ALIGN.LEFT
run = flow.add_run()
run.text = ("find_hits (19 ChEMBL hits)  →  cluster_scaffolds (11 clusters)  →  fold_target (Boltz-1 structure)  "
            "→  dock_compound (docking scores)  →  predict_admet (TIER-1 gate)  →  score_variant_effect (ESM-2 ΔΔG)  "
            "→  generate_pdf_report  ✓")
run.font.size  = Pt(12)
run.font.color.rgb = GREY1


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — AGENT IMPACT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "Agent Impact", "Agents are not peripheral — they are the solution")
footer(sl)

# 2-col comparison
add_rect(sl, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.45), fill=WHITE, line=GREY3)
add_rect(sl, Inches(0.35), Inches(1.35), Inches(5.9), Inches(0.5), fill=GREY3)
add_text(sl, "Without agents", Inches(0.55), Inches(1.4), Inches(5.5), Inches(0.42),
         size=14, bold=True, color=GREY1)

add_rect(sl, Inches(6.55), Inches(1.35), Inches(6.45), Inches(5.45), fill=WHITE, line=REDCLAW_BLUE)
add_rect(sl, Inches(6.55), Inches(1.35), Inches(6.45), Inches(0.5), fill=REDCLAW_BLUE)
add_text(sl, "With Whitespace", Inches(6.75), Inches(1.4), Inches(6.1), Inches(0.42),
         size=14, bold=True, color=WHITE)

without = [
    "Manual cross-referencing of 10+ databases",
    "Days of analyst time per briefing",
    "Fixed pipeline — no adaptive reasoning",
    "No audit trail — not reproducible",
    "One company at a time",
    "Cannot invoke GPU structure prediction",
]
with_ = [
    "Agent autonomously selects & sequences 34 tools",
    "45-minute unattended run",
    "Reasoning adapts per result — different path each time",
    "Full JSON Lines audit trail per session",
    "8 companies in parallel overnight via SLURM sweep",
    "Agent decides when to fold proteins, screen ADMET, dock",
]

txb1 = sl.shapes.add_textbox(Inches(0.55), Inches(1.95), Inches(5.5), Inches(4.6))
tf1 = txb1.text_frame; tf1.word_wrap = True
for w in without:
    bullet(tf1, w, size=13, color=GREY1)

txb2 = sl.shapes.add_textbox(Inches(6.75), Inches(1.95), Inches(6.1), Inches(4.6))
tf2 = txb2.text_frame; tf2.word_wrap = True
for w in with_:
    bullet(tf2, w, size=13, color=DARK_BLUE, bold=True)

# Bottom note
add_rect(sl, Inches(0.35), Inches(6.75), Inches(12.65), Inches(0.4), fill=XL_BLUE)
add_text(sl, "Development was also agent-assisted: Claude Code built proxy retry logic, egress allowlist, audit test suite, guard logic — all during hackathon week.",
         Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.35),
         size=11, color=DARK_BLUE, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WHAT YOU CAN ASK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "What You Can Ask", "Natural-language queries — no forms, no dashboards, no waiting")
footer(sl)

ceo_categories = [
    (DARK_BLUE, "Portfolio & Pipeline", [
        '"Rank our full pipeline by opportunity score. Which 5 assets do we accelerate in the next 90 days?"',
        '"Which Phase 2 assets have the highest probability of Phase 3 success?"',
        '"Show me every RedClaw asset in a competitive vacuum — no rival in the same indication."',
    ]),
    (REDCLAW_BLUE, "Competitive Intelligence", [
        '"Where is AstraZeneca moving faster than us, and what should we do about it?"',
        '"Give me a full briefing on Eli Lilly — where do they beat us, where do we win?"',
        '"Which competitor is most likely to challenge our oncology franchise in 2 years?"',
    ]),
    (LIGHT_BLUE, "White Space & First-Mover", [
        '"Find three therapeutic white spaces where RedClaw has the strongest science and zero major competitor presence."',
        '"Where can RedClaw be first-in-class in the next 5 years based on our biology advantages?"',
    ]),
    (RGBColor(0xFF,0x7D,0x29), "Repurposing", [
        '"Which RedClaw drugs can be repurposed into uncrowded oncology indications — rank by ADMET profile."',
        '"What\'s the single highest-ROI repurposing move RedClaw could make in the next 18 months?"',
    ]),
    (GREEN, "Regulatory & Orphan", [
        '"Which pipeline assets qualify for orphan drug designation and what\'s the revenue uplift?"',
        '"Map the regulatory path to approval for [asset] in the US and EU — key milestones and risks."',
    ]),
    (GREY1, "Full CEO Briefing", [
        '"Portfolio rank · competitive gaps · three white spaces · top 5 trial scores · two biggest threats · single highest-ROI 90-day action. Generate a PDF."',
    ]),
]

# 3-column layout, 2 rows
positions = [
    (Inches(0.3),  Inches(1.35), Inches(4.2)),
    (Inches(4.6),  Inches(1.35), Inches(4.2)),
    (Inches(8.9),  Inches(1.35), Inches(4.1)),
    (Inches(0.3),  Inches(4.1),  Inches(4.2)),
    (Inches(4.6),  Inches(4.1),  Inches(4.2)),
    (Inches(8.9),  Inches(4.1),  Inches(4.1)),
]

for i, (col, cat_title, queries) in enumerate(ceo_categories):
    lx, ly, lw = positions[i]
    lh = Inches(2.55) if ly < Inches(3) else Inches(2.8)
    add_rect(sl, lx, ly, lw, lh, fill=WHITE, line=col)
    add_rect(sl, lx, ly, lw, Inches(0.38), fill=col)
    add_text(sl, cat_title, lx + Inches(0.1), ly + Inches(0.04), lw - Inches(0.15), Inches(0.33),
             size=11, bold=True, color=WHITE)
    txb = sl.shapes.add_textbox(lx + Inches(0.12), ly + Inches(0.45), lw - Inches(0.2), lh - Inches(0.55))
    tf = txb.text_frame; tf.word_wrap = True
    for q in queries:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = q
        run.font.size = Pt(10)
        run.font.color.rgb = GREY1
        run.font.italic = True
        p.space_before = Pt(3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NEAR_WHITE)
header_bar(sl, "Next Steps", "From hackathon prototype to RedClaw R&D platform")
footer(sl)

roadmap = [
    (DARK_BLUE,  "90 days",  "Production hardening",
     ["Eliminate proxy stall (direct API integration)", "ESM-2 GPU streaming fix (clawapi config)", "Grafana dashboard over audit logs"]),
    (REDCLAW_BLUE, "6 months", "Platform expansion",
     ["Integrate with Agenecy as reusable executable skills", "Extend to 12 competitors", "Add real-time ClinicalTrials.gov change monitoring"]),
    (LIGHT_BLUE, "12 months","Full drug discovery loop",
     ["Gap → target → generative chemistry → ADMET → IND briefing in one run", "Patient-level de-identified trial outcomes for better trial success scoring", "Molecule optimisation loop (docking → ADMET → retrosynthesis)"]),
]

x = Inches(0.35)
for col, horizon, title, items in roadmap:
    # horizon chip
    add_rect(sl, x, Inches(1.45), Inches(4.15), Inches(0.45), fill=col)
    add_text(sl, horizon, x + Inches(0.1), Inches(1.47), Inches(3.9), Inches(0.4),
             size=13, bold=True, color=WHITE)
    # card
    add_rect(sl, x, Inches(1.9), Inches(4.15), Inches(4.7), fill=WHITE, line=col)
    add_text(sl, title, x + Inches(0.15), Inches(2.0), Inches(3.85), Inches(0.45),
             size=14, bold=True, color=DARK_BLUE)
    txb = sl.shapes.add_textbox(x + Inches(0.15), Inches(2.5), Inches(3.85), Inches(3.8))
    tf = txb.text_frame; tf.word_wrap = True
    for item in items:
        bullet(tf, item, size=12, color=GREY1)
    x += Inches(4.44)

# Vision statement
add_rect(sl, Inches(0.35), Inches(6.75), Inches(12.65), Inches(0.4), fill=DARK_BLUE)
add_text(sl,
         "Vision: Whitespace becomes the always-on strategic intelligence layer "
         "for every RedClaw R&D decision — not a one-off briefing tool.",
         Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.35),
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BLUE)
add_rect(sl, Inches(8.5), 0, Inches(4.83), SLIDE_H, fill=REDCLAW_BLUE)
add_rect(sl, Inches(8.3), 0, Inches(0.25), SLIDE_H, fill=LIGHT_BLUE)

add_text(sl, "Whitespace",
         Inches(0.6), Inches(2.0), Inches(7.8), Inches(1.2),
         size=48, bold=True, color=WHITE)
add_text(sl, "From fragmented databases to CEO-ready strategy\nin one autonomous agent run.",
         Inches(0.6), Inches(3.15), Inches(7.6), Inches(1.0),
         size=18, color=XL_BLUE, italic=True)

add_text(sl, "34 tools  ·  GPU-accelerated  ·  Audit trail  ·  SLURM-native",
         Inches(0.6), Inches(4.35), Inches(7.6), Inches(0.45),
         size=14, color=GREY3)

add_text(sl, "Thank you",
         Inches(8.85), Inches(3.2), Inches(4.0), Inches(0.8),
         size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "RedClaw CSI Hackathon 2026",
         Inches(0.6), Inches(5.9), Inches(5), Inches(0.4),
         size=12, color=GREY3)
add_text(sl, "April 2026  |  For Internal Use Only",
         Inches(0.6), Inches(6.25), Inches(5), Inches(0.35),
         size=11, color=GREY3)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "/home/sobhn/hk/drug-discovery-agent/reports/Whitespace_Hackathon_2026.pptx"
prs.save(out)
print(f"Saved: {out}")
