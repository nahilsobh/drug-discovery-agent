#!/usr/bin/env python3.11
"""Generate Whitespace hackathon slide deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── RedClaw brand colours ───────────────────────────────────────────────────────
NAVY   = RGBColor(0x02, 0x23, 0x66)   # Dark Blue
BLUE   = RGBColor(0x0B, 0x41, 0xCD)   # RedClaw Blue
LBLUE  = RGBColor(0x14, 0x82, 0xFA)   # Light Blue
XBLUE  = RGBColor(0xBD, 0xE3, 0xFF)   # Extra Light Blue
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x54, 0x4F, 0x4F)
GREEN  = RGBColor(0x00, 0xB4, 0x58)
RED    = RGBColor(0xFF, 0x1F, 0x26)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank)


def rect(slide, x, y, w, h, fill=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def txbox(slide, text, x, y, w, h, size=24, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def bullet_box(slide, items, x, y, w, h, size=18, color=WHITE, title=None, title_size=22):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = color
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.level = item.get("level", 0)
        r = p.add_run()
        r.text = ("• " if item.get("bullet", True) else "") + item["text"]
        r.font.size = Pt(item.get("size", size))
        r.font.bold = item.get("bold", False)
        r.font.color.rgb = item.get("color", color)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 0.08, BLUE)
rect(s, 0, 7.42, 13.33, 0.08, BLUE)
# Accent bar
rect(s, 0, 2.8, 0.12, 2.2, LBLUE)

txbox(s, "WHITESPACE", 0.4, 1.2, 12, 1.2, size=64, bold=True, color=WHITE)
txbox(s, "RedClaw AI Factory — Strategic Discovery Agent", 0.4, 2.5, 12, 0.7, size=26, color=XBLUE)
txbox(s, "34-tool autonomous agent · CEO-ready PDF reports · GPU-accelerated drug discovery",
      0.4, 3.3, 11, 0.6, size=18, color=RGBColor(0xA0, 0xBF, 0xFF))
txbox(s, "CSI Hackathon 2026  ·  April 30", 0.4, 6.8, 8, 0.5, size=14,
      color=RGBColor(0x80, 0xA0, 0xD0))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xFF))
rect(s, 0, 0, 13.33, 1.1, NAVY)
rect(s, 0, 1.1, 13.33, 0.05, BLUE)
txbox(s, "The Problem", 0.4, 0.2, 12, 0.7, size=32, bold=True, color=WHITE)

# Three pain-point cards
cards = [
    ("⏱  Weeks of analyst time",
     "Synthesizing competitor pipelines, trial data, patents, and literature takes weeks per report — by then decisions are already made."),
    ("📊  Fragmented sources",
     "ClinicalTrials.gov, Open Targets, ChEMBL, FDA FAERS, USPTO, Europe PMC — no single system connects them for strategic decisions."),
    ("🔒  No CEO-ready interface",
     "Drug Discovery VPs rely on stale PowerPoint decks. There is no way to ask a live strategic question and get a same-day answer."),
]
for i, (title, body) in enumerate(cards):
    cx = 0.3 + i * 4.35
    rect(s, cx, 1.4, 4.1, 5.6, WHITE)
    # top accent
    rect(s, cx, 1.4, 4.1, 0.12, BLUE)
    txbox(s, title, cx + 0.15, 1.6, 3.8, 0.7, size=17, bold=True, color=NAVY)
    txbox(s, body,  cx + 0.15, 2.4, 3.8, 4.2, size=15, color=GREY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Solution
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x01, 0x18, 0x45))
txbox(s, "The Solution", 0.4, 0.2, 12, 0.7, size=32, bold=True, color=WHITE)

txbox(s, "Ask a strategic question. Get a branded CEO-ready PDF in 30 minutes. Zero human steps.",
      0.4, 1.2, 12.5, 0.7, size=20, bold=True, color=XBLUE)

items = [
    {"text": "34-tool ReAct agent — autonomously decides tool sequence across up to 45 reasoning turns", "bold": True},
    {"text": "Live data: ClinicalTrials.gov · Open Targets · ChEMBL · Europe PMC · FDA FAERS · USPTO · KEGG · Orphanet"},
    {"text": "GPU tools: Boltz-1 protein folding · ESM-2 variant effects · ADMET prediction · scaffold clustering · docking"},
    {"text": "Whitespace Web UI — CEO/VP browser interface: 18 pre-built queries, live log streaming, instant chat, PDF viewer"},
    {"text": "Output: branded RedClaw PDF with ranked recommendations, composite scores, and 90-day action items"},
]
bullet_box(s, items, 0.4, 2.1, 12.5, 4.5, size=17, color=WHITE)

rect(s, 0.4, 6.2, 12.5, 0.9, RGBColor(0x0B, 0x30, 0x80))
txbox(s, "17 branded strategic intelligence PDFs generated autonomously in one 3-hour sweep",
      0.6, 6.25, 12, 0.7, size=18, bold=True, color=GREEN)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Architecture
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xFF))
rect(s, 0, 0, 13.33, 1.1, NAVY)
txbox(s, "Architecture", 0.4, 0.2, 12, 0.7, size=32, bold=True, color=WHITE)

# Flow boxes
flow = [
    ("CEO / VP\nBrowser", BLUE),
    ("Whitespace\nWeb UI\n(FastAPI)", NAVY),
    ("Claude Opus\nReAct Loop\n34 tools", BLUE),
    ("External APIs\n+ GPU", RGBColor(0x00, 0x80, 0x60)),
    ("Branded\nPDF Report", NAVY),
]
arrow_y = 3.5
for i, (label, color) in enumerate(flow):
    bx = 0.3 + i * 2.6
    rect(s, bx, 2.2, 2.2, 1.8, color)
    txbox(s, label, bx + 0.1, 2.35, 2.0, 1.5, size=14, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        txbox(s, "→", bx + 2.2, 2.8, 0.4, 0.6, size=22, bold=True, color=BLUE)

# Tool categories below
cats = [
    ("DISCOVERY", ["find_gaps", "get_biology", "search_redclaw_trials", "find_combinations", "get_pathway_context"]),
    ("COMPETITIVE", ["rank_portfolio", "list_pipeline_assets", "monitor_competitive_signals", "map_regulatory_path", "score_trial_outcome"]),
    ("SCIENCE", ["scan_literature", "scan_arxiv", "find_hits", "query_adverse_events", "find_repurposing_candidates"]),
    ("GPU / GENOMECLAW", ["fold_target", "predict_admet", "score_variant_effect", "cluster_scaffolds", "dock_compound"]),
]
for i, (cat, tools) in enumerate(cats):
    cx = 0.3 + i * 3.25
    rect(s, cx, 4.3, 3.0, 2.9, WHITE)
    rect(s, cx, 4.3, 3.0, 0.35, BLUE)
    txbox(s, cat, cx + 0.08, 4.32, 2.85, 0.3, size=10, bold=True, color=WHITE)
    for j, t in enumerate(tools):
        txbox(s, "· " + t, cx + 0.1, 4.75 + j * 0.42, 2.8, 0.4, size=11, color=GREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Agent Workflow
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x01, 0x18, 0x45))
txbox(s, "How the Agent Works", 0.4, 0.2, 12, 0.7, size=32, bold=True, color=WHITE)

txbox(s, "Competitive Intelligence — RedClaw vs AstraZeneca (example, 13 tool calls, 30 turns)",
      0.4, 1.15, 12.5, 0.5, size=16, color=XBLUE)

steps = [
    ("1", "recall_longterm_memory", "Load prior AZ intelligence from cache"),
    ("2", "list_pipeline_assets", "Enumerate all 59 RedClaw assets"),
    ("3", "rank_portfolio", "Score assets by stage, indication, competitive overlap"),
    ("4", "find_gaps", "Identify indications where AZ is ahead"),
    ("5", "monitor_competitive_signals", "Pull live AZ trial starts, approvals, failures"),
    ("6", "scan_literature", "ArXiv + Europe PMC: AZ mechanism papers (last 18 mo)"),
    ("7", "search_patents", "USPTO freedom-to-operate sweep"),
    ("8", "find_repurposing_candidates", "RedClaw approved drugs for AZ-dominated indications"),
    ("9", "predict_admet ✓ TIER-1", "ADMET gate — mandatory before advancing any compound"),
    ("10", "check_orphan_eligibility", "Rare disease fast-track analysis"),
    ("11", "save_to_cache", "Persist findings to cross-session memory"),
    ("12", "generate_pdf_report", "Branded PDF: ranked table, scores, 90-day action"),
]
cols = 2
per_col = len(steps) // cols
for i, (num, tool, desc) in enumerate(steps):
    col = i // per_col
    row = i % per_col
    bx = 0.4 + col * 6.5
    by = 1.85 + row * 0.46
    rect(s, bx, by, 0.42, 0.36, BLUE)
    txbox(s, num, bx + 0.02, by + 0.02, 0.38, 0.32, size=11, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, tool, bx + 0.5, by + 0.02, 2.3, 0.36, size=12, bold=True, color=XBLUE)
    txbox(s, desc, bx + 2.85, by + 0.02, 3.4, 0.36, size=11, color=RGBColor(0xC0, 0xD0, 0xFF))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Results
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xFF))
rect(s, 0, 0, 13.33, 1.1, NAVY)
txbox(s, "Results — 3-Hour Autonomous Sweep", 0.4, 0.2, 12, 0.7, size=32, bold=True, color=WHITE)

# Stats row
stats = [
    ("17", "PDF reports\ngenerated"),
    ("34", "tools called\nautonomously"),
    ("8", "competitors\nanalyzed"),
    ("0", "human steps\nbetween Q→PDF"),
]
for i, (num, label) in enumerate(stats):
    cx = 0.4 + i * 3.2
    rect(s, cx, 1.3, 2.9, 1.6, NAVY)
    txbox(s, num,   cx + 0.1, 1.35, 2.7, 1.0, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, label, cx + 0.1, 2.2,  2.7, 0.65, size=13, color=XBLUE, align=PP_ALIGN.CENTER)

# Sample finding
rect(s, 0.4, 3.15, 12.5, 1.5, RGBColor(0xE8, 0xF4, 0xEC))
rect(s, 0.4, 3.15, 0.1, 1.5, GREEN)
txbox(s, "Sample Finding — AstraZeneca Repurposing Analysis", 0.6, 3.2, 12, 0.4,
      size=15, bold=True, color=NAVY)
txbox(s, "Capivasertib → Proteus Syndrome  |  Composite score: 0.94  |  AKT1 E17K genotype match  |  "
         "TIER-1 ADMET cleared  |  95% competitive vacuum  |  Orphan-eligible  |  "
         "IND feasible within 18 months  |  7-year market exclusivity",
      0.6, 3.65, 12.1, 0.9, size=14, color=GREY)

# Reports list
txbox(s, "Generated Reports:", 0.4, 4.85, 5, 0.4, size=15, bold=True, color=NAVY)
reports_l = ["CEO Strategic Briefing", "RedClaw vs AstraZeneca", "RedClaw vs Eli Lilly",
             "RedClaw vs Novartis", "RedClaw vs Pfizer", "RedClaw vs Merck",
             "RedClaw vs Bristol-Myers Squibb", "RedClaw vs AbbVie", "RedClaw vs J&J"]
reports_r = ["Repurposing vs AstraZeneca", "Repurposing vs Eli Lilly", "Repurposing vs Novartis",
             "Repurposing vs Pfizer", "Repurposing vs Merck", "Repurposing vs BMS",
             "Repurposing vs AbbVie", "Repurposing vs J&J"]
for i, r in enumerate(reports_l):
    txbox(s, "✓  " + r, 0.4, 5.3 + i * 0.3, 5.5, 0.28, size=12, color=RGBColor(0x00, 0x80, 0x40))
for i, r in enumerate(reports_r):
    txbox(s, "✓  " + r, 6.8, 5.3 + i * 0.3, 5.5, 0.28, size=12, color=RGBColor(0x00, 0x80, 0x40))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Web UI
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x01, 0x18, 0x45))
txbox(s, "Whitespace Web UI", 0.4, 0.2, 12, 0.7, size=32, bold=True, color=WHITE)

# Mock UI diagram
rect(s, 0.3, 1.2, 12.7, 5.9, RGBColor(0xF0, 0xF4, 0xF8))  # bg
rect(s, 0.3, 1.2, 12.7, 0.5, RGBColor(0x00, 0x30, 0x87))  # header
txbox(s, "WHITESPACE  ·  RedClaw AI Factory · Strategic Discovery       ✓ 17 done",
      0.5, 1.25, 12.3, 0.4, size=12, bold=True, color=WHITE)

# Left panel
rect(s, 0.3, 1.7, 2.8, 5.4, WHITE)
txbox(s, "STRATEGIC QUERIES", 0.4, 1.8, 2.6, 0.3, size=9, bold=True, color=GREY)
for i, q in enumerate(["Executive", "Full CEO Briefing", "Full Pipeline Suite",
                        "Competitive Intelligence", "vs AstraZeneca", "vs Eli Lilly",
                        "Drug Repurposing", "vs AstraZeneca "]):
    bold = q in ("Executive", "Competitive Intelligence", "Drug Repurposing")
    col = NAVY if bold else GREY
    sz = 10 if bold else 11
    txbox(s, q, 0.4 + (0 if bold else 0.15), 2.15 + i * 0.38, 2.6, 0.35,
          size=sz, bold=bold, color=col)
rect(s, 0.35, 5.55, 2.7, 1.3, RGBColor(0xF0, 0xF4, 0xF8))
txbox(s, "Ask anything…", 0.45, 5.6, 2.5, 0.5, size=10, color=GREY)
rect(s, 0.45, 6.15, 1.2, 0.35, BLUE)
txbox(s, "Ask", 0.45, 6.17, 1.2, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 1.75, 6.15, 1.2, 0.35, RGBColor(0xE0, 0xE8, 0xF8))
txbox(s, "Run as Job", 1.75, 6.17, 1.2, 0.3, size=10, color=NAVY, align=PP_ALIGN.CENTER)

# Center panel
rect(s, 3.1, 1.7, 5.5, 5.4, RGBColor(0xF8, 0xFA, 0xFC))
txbox(s, "▶ Submitting: RedClaw vs AstraZeneca…", 3.2, 1.85, 5.3, 0.35, size=11, color=GREY)
txbox(s, "✓ Job 34031234 queued — waiting for log…", 3.2, 2.2, 5.3, 0.3, size=11, color=GREY)
for i, line in enumerate([
    "══ Turn 1/30  tools called: 0 ══",
    "│  list_pipeline_assets  →  59 RedClaw assets",
    "══ Turn 4/30  tools called: 3 ══",
    "│  find_gaps  →  12 gaps identified",
    "══ Turn 9/30  tools called: 7 ══",
    "│  predict_admet  →  TIER-1 cleared",
    "══ Turn 13/30  tools called: 11 ══",
    "│  generate_pdf_report  →  PDF saved",
    "[result] SUCCESS — PDF generated",
]):
    col = GREEN if "[result] SUCCESS" in line else (BLUE if "Turn" in line else GREY)
    txbox(s, line, 3.2, 2.6 + i * 0.38, 5.3, 0.35, size=10, color=col)
rect(s, 3.2, 6.05, 2.5, 0.38, GREEN)
txbox(s, "Open PDF: RedClaw_vs_astrazeneca.pdf", 3.25, 6.07, 2.4, 0.32,
      size=10, bold=True, color=WHITE)

# Right panel
rect(s, 8.6, 1.7, 4.4, 5.4, WHITE)
txbox(s, "REPORTS", 8.7, 1.8, 4.2, 0.3, size=9, bold=True, color=GREY)
rpts = ["CEO Briefing 20260430", "RedClaw vs abbvie 20260430",
        "RedClaw vs jnj 20260430", "Repurposing abbvie 20260430",
        "Repurposing bms 20260430", "RedClaw vs bms 20260430"]
for i, r in enumerate(rpts):
    rect(s, 8.6, 2.15 + i * 0.38, 4.4, 0.35, RGBColor(0xF0, 0xF4, 0xF8) if i % 2 == 0 else WHITE)
    txbox(s, r, 8.7, 2.17 + i * 0.38, 4.2, 0.3, size=10, color=NAVY)
rect(s, 8.6, 4.45, 4.4, 2.55, RGBColor(0xE8, 0xF0, 0xFB))
txbox(s, "[ PDF Viewer ]", 8.6, 5.5, 4.4, 0.6, size=16, color=GREY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Agent Impact
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xFF))
rect(s, 0, 0, 13.33, 1.1, NAVY)
txbox(s, "Agent Impact", 0.4, 0.2, 12, 0.7, size=32, bold=True, color=WHITE)

rect(s, 0.4, 1.25, 12.5, 0.8, NAVY)
txbox(s, '"This is not a chatbot with tool access. It is an autonomous pharmaceutical research analyst."',
      0.6, 1.3, 12.1, 0.7, size=17, color=WHITE)

items = [
    ("Fully autonomous reasoning",
     "Agent decides which of 34 tools to call, in what order, across 45 turns — no scripting, no human checkpoints"),
    ("Self-enforcing quality gates",
     "ADMET is a mandatory gate: agent refuses to advance any compound without TIER-1 clearance, even when prompted"),
    ("Self-correction",
     "Proxy stall detection + retry loop: when claude -p returns prose instead of JSON, agent reformats and continues"),
    ("Cross-session memory",
     "Findings saved to agent_longterm_memory.json — each run builds on prior intelligence, compounding over time"),
    ("Agent-built by agents",
     "The entire codebase (34 tools, proxy, web UI, SLURM scripts, 288 tests) was developed using Claude Code as agentic coding assistant"),
]
for i, (title, body) in enumerate(items):
    by = 2.25 + i * 0.95
    rect(s, 0.4, by, 0.08, 0.75, BLUE)
    txbox(s, title, 0.65, by, 4.5, 0.38, size=15, bold=True, color=NAVY)
    txbox(s, body,  0.65, by + 0.38, 12.0, 0.52, size=13, color=GREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Next Steps
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x01, 0x18, 0x45))
txbox(s, "Impact & Next Steps", 0.4, 0.2, 12, 0.7, size=32, bold=True, color=WHITE)

now = [
    "17 CEO-ready PDFs generated in one 3-hour autonomous sweep",
    "Full competitive intelligence vs 8 top pharma companies",
    "Drug repurposing analysis: ADMET-gated, orphan-eligible candidates identified",
    "Capivasertib → Proteus Syndrome: score 0.94, IND-ready in 18 months",
    "Whitespace Web UI live — accessible to non-technical stakeholders today",
]
next_steps = [
    "AWS deployment — remove sHPC dependency, share with all colleagues",
    "RedClaw internal pipeline integration — live asset updates replace static JSON",
    "Expand to 34-company sweep on demand — full market landscape in hours",
    "Connect to Agenecy — make tools reusable across RedClaw R&D organizations",
    "Regulatory filing assistant — map each gap to fastest approval pathway",
]

rect(s, 0.3, 1.2, 6.1, 6.0, RGBColor(0x01, 0x18, 0x55))
rect(s, 6.9, 1.2, 6.1, 6.0, RGBColor(0x01, 0x18, 0x55))

txbox(s, "Delivered this hackathon", 0.5, 1.3, 5.7, 0.5, size=17, bold=True, color=GREEN)
for i, item in enumerate(now):
    txbox(s, "✓  " + item, 0.5, 1.9 + i * 0.82, 5.7, 0.75, size=13, color=WHITE)

txbox(s, "Next steps", 7.1, 1.3, 5.7, 0.5, size=17, bold=True, color=XBLUE)
for i, item in enumerate(next_steps):
    txbox(s, "→  " + item, 7.1, 1.9 + i * 0.82, 5.7, 0.75, size=13, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Closing
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 0.08, BLUE)
rect(s, 0, 7.42, 13.33, 0.08, BLUE)
rect(s, 0, 0, 0.12, 7.5, LBLUE)

txbox(s, "WHITESPACE", 0.4, 1.8, 12, 1.0, size=64, bold=True, color=WHITE)
txbox(s, "From question to CEO-ready strategic intelligence in 30 minutes.", 0.4, 3.0, 12, 0.6,
      size=24, color=XBLUE)
txbox(s, "34 tools  ·  17 reports  ·  8 competitors  ·  0 human steps", 0.4, 3.75, 12, 0.5,
      size=18, color=RGBColor(0xA0, 0xBF, 0xFF))
txbox(s, "CSI Hackathon 2026", 0.4, 6.7, 8, 0.5, size=14,
      color=RGBColor(0x80, 0xA0, 0xD0))


out = "/home/sobhn/hk/drug-discovery-agent/Whitespace_Hackathon_2026.pptx"
prs.save(out)
print(f"Saved: {out}")
